
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import Response
from fastapi.middleware.cors import CORSMiddleware
import tempfile
import os

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def health():
    return {"status": "ok"}


@app.post("/api/convert")
async def convert(file: UploadFile = File(...), outputFormat: str = Form(...)):
    content = await file.read()

    # Save input file
    suffix = "." + file.filename.split(".")[-1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_in:
        tmp_in.write(content)
        input_path = tmp_in.name

    output_path = input_path.rsplit(".", 1)[0] + f".{outputFormat}"

    try:

        # ─── PDF to DOCX ───
        if outputFormat == "docx":
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        # ─── PDF to JPG ───
        elif outputFormat == "jpg":
            from pdf2image import convert_from_path
            import zipfile
            images = convert_from_path(input_path, dpi=150)
            if len(images) == 1:
                images[0].save(output_path, "JPEG")
                media_type = "image/jpeg"
            else:
                # Multiple pages → zip
                zip_path = input_path + ".zip"
                with zipfile.ZipFile(zip_path, "w") as zf:
                    for i, img in enumerate(images):
                        pg_path = input_path + f"_page{i+1}.jpg"
                        img.save(pg_path, "JPEG")
                        zf.write(pg_path, f"page{i+1}.jpg")
                        os.unlink(pg_path)
                output_path = zip_path
                outputFormat = "zip"
                media_type = "application/zip"

        # ─── PDF to XLSX ───
        elif outputFormat == "xlsx":
            import camelot
            import openpyxl
            tables = camelot.read_pdf(input_path, pages="all")
            wb = openpyxl.Workbook()
            ws = wb.active
            for table in tables:
                for row in table.df.values.tolist():
                    ws.append(row)
                ws.append([])
            wb.save(output_path)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

        # ─── PDF to PPTX ───
        elif outputFormat == "pptx":
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches
            images = convert_from_path(input_path, dpi=150)
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            for img in images:
                slide = prs.slides.add_slide(blank_layout)
                img_path = input_path + "_slide.jpg"
                img.save(img_path, "JPEG")
                slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                         prs.slide_width, prs.slide_height)
                os.unlink(img_path)
            prs.save(output_path)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        # ─── DOCX to PDF ───
        elif outputFormat == "pdf" and suffix == ".docx":
            import subprocess
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", os.path.dirname(input_path), input_path], check=True)
            output_path = input_path.replace(".docx", ".pdf")
            media_type = "application/pdf"

        # ─── XLSX to PDF ───
        elif outputFormat == "pdf" and suffix == ".xlsx":
            import subprocess
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", os.path.dirname(input_path), input_path], check=True)
            output_path = input_path.replace(".xlsx", ".pdf")
            media_type = "application/pdf"

        # ─── PPTX to PDF ───
        elif outputFormat == "pdf" and suffix == ".pptx":
            import subprocess
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", os.path.dirname(input_path), input_path], check=True)
            output_path = input_path.replace(".pptx", ".pdf")
            media_type = "application/pdf"

        # ─── JPG/PNG to PDF ───
        elif outputFormat == "pdf" and suffix in (".jpg", ".jpeg", ".png"):
            from PIL import Image
            img = Image.open(input_path).convert("RGB")
            img.save(output_path, "PDF")
            media_type = "application/pdf"

        else:
            return Response(content=b"Unsupported conversion", status_code=400)

        with open(output_path, "rb") as f:
            result = f.read()

        return Response(
            content=result,
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename=converted.{outputFormat}"}
        )

    finally:
        if os.path.exists(input_path):
            os.unlink(input_path)
        if os.path.exists(output_path):
            os.unlink(output_path)
